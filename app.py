import os
import uuid
import base64
import threading
import subprocess
import re
import io
import zipfile
import requests
import time
from flask import Flask, render_template, request, jsonify, send_file, session, redirect, url_for
from pptx import Presentation
from pdf2image import convert_from_path
from functools import wraps
from PIL import Image, ImageFilter

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 150 * 1024 * 1024  # 150MB
app.secret_key = os.environ.get('SECRET_KEY', 'narrador-ai-secret-2025')

JOBS_DIR = '/tmp/jobs'
os.makedirs(JOBS_DIR, exist_ok=True)

jobs = {}
jobs_lock = threading.Lock()

# --- Env vars ---
# Suporta 'chave-api', 'NVIDIA_API_KEY' e 'GROQ_API_KEY' (para manter compatibilidade)
ENV_API_KEY = (os.environ.get('chave-api', '') or os.environ.get('NVIDIA_API_KEY', '') or os.environ.get('GROQ_API_KEY', '')).strip()

API_BASE_URL = os.environ.get('API_BASE_URL', 'https://api.groq.com/openai/v1').strip().rstrip('/')
DEFAULT_MODEL = os.environ.get('DEFAULT_MODEL', 'openai/gpt-oss-120b').strip()

# Suporte a múltiplos usuários via APP_USUARIO_01/APP_SENHA_01 ... APP_USUARIO_99/APP_SENHA_99
# Também aceita o formato legado APP_USUARIO/APP_SENHA como usuário único
USERS = {}
for i in range(1, 100):
    suffix = f'{i:02d}'
    u = os.environ.get(f'APP_USUARIO_{suffix}', '').strip()
    s = os.environ.get(f'APP_SENHA_{suffix}', '').strip()
    if u and s:
        USERS[u] = s

# Fallback legado: APP_USUARIO / APP_SENHA
_leg_u = os.environ.get('APP_USUARIO', '').strip()
_leg_s = os.environ.get('APP_SENHA', '').strip()
if _leg_u and _leg_s and _leg_u not in USERS:
    USERS[_leg_u] = _leg_s

# Suporte simples: USUARIO + SENHA (mais fácil de configurar no EasyPanel)
_simple_u = os.environ.get('USUARIO', '').strip()
_simple_s = os.environ.get('SENHA', '').strip()
if _simple_u and _simple_s and _simple_u not in USERS:
    USERS[_simple_u] = _simple_s

AUTH_ENABLED = bool(USERS)

# Marcas/watermarks conhecidas para remover do texto gerado
WATERMARK_PATTERNS = [
    r'\bNotebookLM\b',
    r'\bGoogle\s+NotebookLM\b',
    r'\[watermark\]',
    r'\[marca\s+d.?água\]',
]

# --- Auth helper ---
def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if AUTH_ENABLED and not session.get('logged_in'):
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated


def clean_watermarks(text: str) -> str:
    """Remove menções a watermarks/marcas conhecidas do texto gerado."""
    for pattern in WATERMARK_PATTERNS:
        text = re.sub(pattern, '', text, flags=re.IGNORECASE)
    # Remove espaços duplicados resultantes
    text = re.sub(r'  +', ' ', text).strip()
    return text


# Termos que identificam shapes de watermark visual nos slides
WATERMARK_SHAPE_TERMS = [
    'notebooklm', 'notebook lm',
]


def _element_contains_watermark(element) -> bool:
    """Busca recursivamente no XML se algum nó de texto contém termo de watermark."""
    try:
        # Serializa todo o XML do elemento e busca pelos termos
        xml_str = element.xml.lower()
        return any(term in xml_str for term in WATERMARK_SHAPE_TERMS)
    except Exception:
        return False


def _is_bottom_right_shape(shape, slide_width, slide_height) -> bool:
    """Retorna True se o shape está no canto inferior direito (onde o NotebookLM fica)."""
    try:
        # Considera bottom-right se estiver nos 30% direitos e 20% inferiores
        return (
            shape.left > slide_width * 0.65 and
            shape.top  > slide_height * 0.75
        )
    except Exception:
        return False


def _remove_shapes_from_collection(shapes, slide_width=None, slide_height=None) -> int:
    """Remove shapes de watermark de uma coleção; retorna quantidade removida."""
    removed = 0
    to_remove = []
    for shape in shapes:
        # Verifica pelo XML completo (pega grupos, imagens com texto alternativo, etc.)
        if _element_contains_watermark(shape._element):
            to_remove.append(shape)
        # Fallback: shape pequeno no canto inferior direito
        elif slide_width and slide_height and _is_bottom_right_shape(shape, slide_width, slide_height):
            # Só remove se for pequeno (não é um bloco de conteúdo real)
            try:
                area = shape.width * shape.height
                slide_area = slide_width * slide_height
                if area < slide_area * 0.05:   # menor que 5% da área do slide
                    to_remove.append(shape)
            except Exception:
                pass

    for shape in to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
            removed += 1
        except Exception:
            pass
    return removed


def remove_watermark_shapes(prs) -> int:
    """
    Remove shapes de watermark visual (ex: logo NotebookLM) da apresentação.
    Verifica slides individuais, layouts e o slide master.
    Retorna o total de shapes removidos.
    """
    total = 0

    # Dimensões do slide para heurística de posição
    slide_width  = prs.slide_width
    slide_height = prs.slide_height

    # 1. Slides individuais
    for slide in prs.slides:
        total += _remove_shapes_from_collection(slide.shapes, slide_width, slide_height)

    # 2. Layouts de slide
    for layout in prs.slide_layouts:
        total += _remove_shapes_from_collection(layout.shapes, slide_width, slide_height)

    # 3. Slide master
    try:
        total += _remove_shapes_from_collection(prs.slide_master.shapes, slide_width, slide_height)
    except Exception:
        pass

    return total


def update_job(job_id, **kwargs):
    with jobs_lock:
        if job_id in jobs:
            jobs[job_id].update(kwargs)


def _paint_watermark_region(img_bytes: bytes, img_format: str) -> bytes:
    """
    Mesma técnica do notebooklmremover.com:
    Abre a imagem do slide, detecta a região do watermark NotebookLM
    (canto inferior direito, proporcional ao tamanho), amostra a cor
    de fundo ao redor e pinta um retângulo cobrindo o logo.
    """
    img = Image.open(io.BytesIO(img_bytes)).convert('RGB')
    w, h = img.size

    # Coordenadas proporcionais ao tamanho (conforme o código deles)
    # Modo landscape (apresentações normais)
    scale = w / 2880
    x = int(2640 * scale)   # left do retângulo
    y = int(1550 * scale)   # top do retângulo
    rw = int(240 * scale)   # largura
    rh = int(50 * scale)    # altura

    # Garante que está dentro dos limites
    x  = max(0, min(x, w - rw))
    y  = max(0, min(y, h - rh))
    rw = min(rw, w - x)
    rh = min(rh, h - y)

    if rw <= 0 or rh <= 0:
        return img_bytes

    # Amostra a cor média da linha logo acima do watermark (igual ao site)
    sample_y = max(0, y - 5)
    sample_region = img.crop((x, sample_y, x + rw, sample_y + 2))
    pixels = list(sample_region.getdata())
    if pixels:
        r_avg = sum(p[0] for p in pixels) // len(pixels)
        g_avg = sum(p[1] for p in pixels) // len(pixels)
        b_avg = sum(p[2] for p in pixels) // len(pixels)
        bg_color = (r_avg, g_avg, b_avg)
    else:
        bg_color = (255, 255, 255)

    # Pinta o retângulo sobre o watermark
    from PIL import ImageDraw
    draw = ImageDraw.Draw(img)
    # Borda suavizada: rect ligeiramente maior
    draw.rectangle([x - 2, y - 2, x + rw + 2, y + rh + 2], fill=bg_color)

    out = io.BytesIO()
    fmt = 'JPEG' if img_format.lower() in ('jpg', 'jpeg') else 'PNG'
    img.save(out, format=fmt, quality=95)
    return out.getvalue()


def remove_watermark_from_images(pptx_path: str) -> str:
    """
    Abre o PPTX como ZIP, processa cada imagem de slide para
    cobrir o watermark NotebookLM, e salva um novo arquivo PPTX.
    Retorna o caminho do arquivo modificado.
    """
    out_path = pptx_path.replace('.pptx', '_wm_clean.pptx')
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(out_path, 'w', compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                name_lower = item.filename.lower()
                # Processa imagens de mídia dos slides
                if (name_lower.startswith('ppt/media/image') and
                        any(name_lower.endswith(ext) for ext in ('.png', '.jpg', '.jpeg'))):
                    ext = name_lower.rsplit('.', 1)[-1]
                    try:
                        data = _paint_watermark_region(data, ext)
                    except Exception:
                        pass  # Se falhar, usa a imagem original
                zout.writestr(item, data)
    return out_path


def process_presentation(job_id, pptx_path, api_key, model, total_slides):
    job_dir = os.path.join(JOBS_DIR, job_id)

    try:
        update_job(job_id, status='converting', message='Removendo marca d\'água das imagens...')

        # Passo 1: Remove watermark das imagens dentro do PPTX (técnica do notebooklmremover.com)
        try:
            clean_pptx_path = remove_watermark_from_images(pptx_path)
        except Exception:
            clean_pptx_path = pptx_path  # fallback: usa o original

        update_job(job_id, message='Convertendo apresentação para imagens...')

        # Passo 2: Convert PPTX limpo → PDF via LibreOffice
        subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', job_dir, clean_pptx_path],
            capture_output=True, text=True, timeout=180
        )

        pdf_name = os.path.splitext(os.path.basename(clean_pptx_path))[0] + '.pdf'
        pdf_path = os.path.join(job_dir, pdf_name)


        if not os.path.exists(pdf_path):
            update_job(job_id, status='error', message='Falha ao converter o arquivo. Verifique se é um .pptx válido.')
            return

        update_job(job_id, message='Renderizando slides...')
        images = convert_from_path(pdf_path, dpi=150)

        prs = Presentation(clean_pptx_path)
        slide_count = min(len(prs.slides), len(images))

        # Remove logos/watermarks visuais dos slides antes de processar
        wm_removed = remove_watermark_shapes(prs)
        if wm_removed:
            update_job(job_id, message=f'Removendo {wm_removed} marca(s) d\'água dos slides...')

        for i in range(slide_count):
            slide = prs.slides[i]
            image = images[i]

            update_job(
                job_id,
                status='processing',
                current_slide=i + 1,
                message=f'Gerando notas para o slide {i + 1} de {slide_count}...'
            )

            img_path = os.path.join(job_dir, f'slide_{i}.png')
            image.save(img_path, 'PNG')

            with open(img_path, 'rb') as f:
                image_data = base64.b64encode(f.read()).decode()

            # Extrai texto do slide via python-pptx (para modelos sem suporte a visão)
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                slide_texts.append(text)
            slide_text_content = '\n'.join(slide_texts) if slide_texts else '(slide sem texto visível)'

            # Instrução de posição do slide (primeiro / meio / último)
            is_first = (i == 0)
            is_last  = (i == slide_count - 1)

            if is_first:
                position_instruction = (
                    "- Este é o PRIMEIRO slide. Inicie o assunto de forma direta ou com uma introdução ao tema. "
                    "NÃO gere saudações iniciais como 'bom dia', 'boa tarde', 'boa noite', 'olá' ou 'sejam bem-vindos'."
                )
            elif is_last:
                position_instruction = (
                    "- Este é o ÚLTIMO slide. Encerre as notas com uma despedida natural e definitiva, "
                    "como 'Ficamos por aqui.' ou 'É isso por hoje.' "
                    "NÃO use frases como 'até a próxima aula', 'nos vemos em breve' ou qualquer variação "
                    "que pressuponha que haverá um próximo encontro — pode ser a última aula da disciplina."
                )
            else:
                position_instruction = (
                    "- Este é um slide intermediário. NÃO faça saudações, boas-vindas ou despedidas. "
                    "Vá direto ao conteúdo do slide."
                )

            # Lê configurações de texto das variáveis de ambiente
            multiplicador   = os.environ.get('MULTIPLICADOR_TEXTO', '2')
            limite_palavras = int(os.environ.get('LIMITE_PALAVRAS', '120'))
            max_tok         = min(int(limite_palavras * 1.8), 1024)

            prompt_text = f"""Você é um especialista em criação de roteiros para apresentações profissionais.

Analise o conteúdo do slide {i + 1} de {slide_count} e crie notas de narrador em português do Brasil.

CONTEÚDO DO SLIDE:
\"\"\"
{slide_text_content}
\"\"\"

Regras OBRIGATÓRIAS:
- Escreva aproximadamente {multiplicador}x o volume de palavras visíveis no próprio slide.
- Use no MÁXIMO {limite_palavras} palavras no total. Se o slide tem poucas palavras, as notas devem ser curtas também.
- Explique e desenvolva o conteúdo do slide de forma natural para fala, sem extrapolar muito além do que está mostrado.
- Inclua frases de transição suaves e naturais.
- Mantenha tom profissional e acessível.
- NUNCA use saudações temporais (como "bom dia", "boa tarde", "boa noite" ou "olá"). O conteúdo será usado por alunos EAD que poderão assistir em qualquer horário.
- IGNORE completamente quaisquer marcas d'água, logos, rodapés ou watermarks visíveis no slide (ex: NotebookLM, Google, etc.) — não os mencione de forma alguma.
{position_instruction}
- Retorne SOMENTE o texto corrido das notas, sem títulos, marcadores ou qualquer formatação."""

            try:
                headers = {
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json"
                }

                # Tenta primeiro com imagem (para modelos com suporte a visão)
                payload = {
                    "model": model,
                    "messages": [{
                        "role": "user",
                        "content": [
                            {
                                "type": "image_url",
                                "image_url": {"url": f"data:image/png;base64,{image_data}"}
                            },
                            {
                                "type": "text",
                                "text": prompt_text
                            }
                        ]
                    }],
                    "max_tokens": max_tok,
                    "temperature": 1,
                    "top_p": 0.95
                }

                resp_api = requests.post(f"{API_BASE_URL}/chat/completions", headers=headers, json=payload, timeout=60)

                # Se o modelo não suporta visão (400), tenta modo texto puro
                if resp_api.status_code == 400:
                    payload = {
                        "model": model,
                        "messages": [{
                            "role": "user",
                            "content": prompt_text
                        }],
                        "max_tokens": max_tok,
                        "temperature": 1,
                        "top_p": 0.95
                    }
                    resp_api = requests.post(f"{API_BASE_URL}/chat/completions", headers=headers, json=payload, timeout=60)

                resp_api.raise_for_status()
                notes_text = resp_api.json()['choices'][0]['message']['content'].strip()
                notes_text = clean_watermarks(notes_text)

            except Exception as e:
                notes_text = f"[Erro ao gerar notas para este slide: {str(e)}]"

            # Insert notes into the slide
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.text = notes_text
            
            # Respeitar limite de requisições (ajustável se necessário)
            time.sleep(2)

        output_path = os.path.join(job_dir, 'output.pptx')
        prs.save(output_path)

        update_job(
            job_id,
            status='done',
            message='Apresentação pronta!',
            current_slide=slide_count,
            output_path=output_path
        )

    except Exception as e:
        update_job(job_id, status='error', message=f'Erro inesperado: {str(e)}')



# ---------------------------------------------------------------------------
# Routes — Auth
# ---------------------------------------------------------------------------

@app.route('/login', methods=['GET', 'POST'])
def login():
    if not AUTH_ENABLED:
        return redirect(url_for('index'))

    error = None
    if request.method == 'POST':
        usuario = request.form.get('usuario', '').strip()
        senha   = request.form.get('senha', '').strip()
        if USERS.get(usuario) == senha:
            session['logged_in'] = True
            session['usuario'] = usuario
            return redirect(url_for('index'))
        error = 'Usuário ou senha incorretos.'

    return render_template('login.html', error=error)


@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('login'))


# ---------------------------------------------------------------------------
# Routes — Main
# ---------------------------------------------------------------------------

@app.route('/')
@login_required
def index():
    # Informa ao template se a API key está pré-configurada no ambiente
    has_env_key = bool(ENV_API_KEY)
    auth_enabled = AUTH_ENABLED
    return render_template('index.html', has_env_key=has_env_key, auth_enabled=auth_enabled, default_model=DEFAULT_MODEL)


@app.route('/models')
@login_required
def get_models():
    """Busca modelos disponíveis na API do provedor."""
    api_key = ENV_API_KEY or request.headers.get('X-Api-Key', '').strip()
    if not api_key:
        return jsonify({'error': 'API key não disponível.'}), 400

    try:
        resp = requests.get(
            f'{API_BASE_URL}/models',
            headers={'Authorization': f'Bearer {api_key}'},
            timeout=10
        )
        resp.raise_for_status()
        data = resp.json()

        # Variável de ambiente para filtrar modelos específicos (separados por vírgula)
        allowed_models_env = os.environ.get('MODELOS_PERMITIDOS', '').strip()
        allowed_models = [m.strip().lower() for m in allowed_models_env.split(',') if m.strip()]

        # Retornamos os modelos disponíveis
        vision_models = []
        for m in data.get('data', []):
            model_id = m.get('id', '')
            model_id_lower = model_id.lower()
            
            if allowed_models:
                # Se o usuário definiu modelos permitidos, respeitar essa lista
                if any(allowed in model_id_lower for allowed in allowed_models):
                    vision_models.append({
                        'id': model_id,
                        'label': model_id
                    })
            else:
                # Filtramos apenas modelos que comprovadamente suportam visão
                # (Modelos de texto puros retornariam erro ao receber imagem)
                is_vision = (
                    'vision' in model_id_lower or
                    'vl' in model_id_lower or
                    'vila' in model_id_lower or
                    'neva' in model_id_lower or
                    'pixtral' in model_id_lower
                )
                if is_vision:
                    vision_models.append({
                        'id': model_id,
                        'label': model_id
                    })
        
        # Ordena alfabeticamente
        vision_models = sorted(vision_models, key=lambda x: x['label'])

        # Se a API não retornou nenhum (fallback para lista conhecida)
        if not vision_models:
            vision_models = [
                {'id': DEFAULT_MODEL, 'label': DEFAULT_MODEL},
            ]

        return jsonify({'models': vision_models})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/upload', methods=['POST'])
@login_required
def upload():
    # Usa a API key do ambiente se disponível, senão pega do form
    api_key = ENV_API_KEY or request.form.get('api_key', '').strip()
    model   = request.form.get('model', DEFAULT_MODEL).strip()
    file    = request.files.get('file')

    if not api_key:
        return jsonify({'error': 'A API key do provedor é obrigatória.'}), 400
    if not file:
        return jsonify({'error': 'Nenhum arquivo enviado.'}), 400
    if not file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'Apenas arquivos .pptx são suportados.'}), 400

    job_id  = str(uuid.uuid4())
    job_dir = os.path.join(JOBS_DIR, job_id)
    os.makedirs(job_dir, exist_ok=True)

    # Gera o nome de download preservando o nome original + sufixo _notas
    original_stem    = os.path.splitext(file.filename)[0]
    download_name    = f"{original_stem}_notas.pptx"

    pptx_path = os.path.join(job_dir, 'presentation.pptx')
    file.save(pptx_path)

    try:
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
    except Exception:
        return jsonify({'error': 'Arquivo inválido ou corrompido.'}), 400

    with jobs_lock:
        jobs[job_id] = {
            'status': 'queued',
            'message': 'Na fila...',
            'current_slide': 0,
            'total_slides': total_slides,
            'output_path': None,
            'download_name': download_name
        }

    thread = threading.Thread(
        target=process_presentation,
        args=(job_id, pptx_path, api_key, model, total_slides)
    )
    thread.daemon = True
    thread.start()

    return jsonify({'job_id': job_id, 'total_slides': total_slides})


@app.route('/status/<job_id>')
@login_required
def status(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return jsonify({'error': 'Job não encontrado.'}), 404
    return jsonify({
        'status': job['status'],
        'message': job['message'],
        'current_slide': job['current_slide'],
        'total_slides': job['total_slides']
    })


@app.route('/download/<job_id>')
@login_required
def download(job_id):
    with jobs_lock:
        job = jobs.get(job_id)
    if not job:
        return jsonify({'error': 'Job não encontrado.'}), 404
    if job['status'] != 'done':
        return jsonify({'error': 'Arquivo ainda não está pronto.'}), 400

    output_path = job.get('output_path')
    if not output_path or not os.path.exists(output_path):
        return jsonify({'error': 'Arquivo de saída não encontrado.'}), 404

    # Usa o nome original do arquivo + sufixo _notas
    download_name = job.get('download_name', 'apresentacao_notas.pptx')

    return send_file(
        output_path,
        as_attachment=True,
        download_name=download_name,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )


@app.route("/health")
def health():
    return jsonify({"status": "ok"}), 200


if __name__ == '__main__':
    app.run(host="0.0.0.0", port=5000, debug=False, threaded=True)
